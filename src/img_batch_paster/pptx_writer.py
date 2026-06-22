from __future__ import annotations

import copy
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Cm

from .config import Config
from .grouper import GroupedImages


_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _remap_rids_in_tree(elem, rid_map: dict[str, str]) -> None:
    """Walk an XML subtree and rewrite r:* attributes whose value appears in rid_map."""
    for el in elem.iter():
        for attr, val in list(el.attrib.items()):
            if attr.startswith(_R_NS) and val in rid_map:
                el.attrib[attr] = rid_map[val]


def _duplicate_slide(prs, source_slide):
    """Duplicate an existing slide: copy shapes and their relationships, remapping rIds.

    python-pptx assigns its own rIds when we add relationships to the new slide, but the
    deep-copied shape XML still references the source's rIds. We build a src→dst rId map
    and rewrite the copied XML so all references resolve. Otherwise PowerPoint complains
    about needing repair (dangling r:embed / r:link references).
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    # Drop default placeholders coming from the layout
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    src_rels = source_slide.part.rels
    dst_rels = new_slide.part.rels

    # Build rId map by ensuring each source relationship exists in the new slide
    rid_map: dict[str, str] = {}
    for rel in src_rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = dst_rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = dst_rels.get_or_add(rel.reltype, rel.target_part)
        if new_rid != rel.rId:
            rid_map[rel.rId] = new_rid

    # Deep-copy shapes and remap their rId references
    for shp in source_slide.shapes:
        new_el = copy.deepcopy(shp._element)
        if rid_map:
            _remap_rids_in_tree(new_el, rid_map)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


@dataclass
class Placement:
    path: Path | None
    x_cm: float
    y_cm: float
    w_cm: float
    h_cm: float
    text: str | None = None       # 若設了 text，就會插入文字方塊而不是圖片
    font_pt: float = 18.0
    bold: bool = True
    align: str = "center"         # "left" | "center" | "right"
    row_idx: int | None = None    # SN 文字所屬的列 index (給「寫入表格 cell」用)
    crop: dict | None = None      # per-placement crop override (覆寫全域 crop)


def _find_table(slide):
    for shp in slide.shapes:
        try:
            if shp.has_table:
                return shp.table
        except Exception:
            continue
    return None


def _write_sn_into_table(slide, placements, sn_col: int, sn_row_start: int) -> bool:
    """把含 text + row_idx 的 placement 寫進 slide 第一個表格的儲存格。
    回傳 True 表示有寫入 (table 存在)；False 則呼叫端應 fallback 用文字方塊。"""
    from pptx.util import Pt
    table = _find_table(slide)
    if table is None:
        return False
    nrows, ncols = len(table.rows), len(table.columns)
    for pl in placements:
        if pl.text is None or pl.row_idx is None:
            continue
        tr = sn_row_start + pl.row_idx
        if tr < 0 or tr >= nrows or sn_col < 0 or sn_col >= ncols:
            continue
        cell = table.cell(tr, sn_col)
        cell.text = str(pl.text)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(pl.font_pt)
                r.font.bold = pl.bold
    return True


def _add_placements_to_slide(slide, placements: list[Placement],
                             sn_in_cell: bool = False,
                             sn_col: int = 0, sn_row_start: int = 1,
                             crop: dict | None = None,
                             tmp_dir: Path | None = None) -> None:
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Pt
    from .xlsx_writer import _apply_crop
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    # SN 寫入表格 cell（成功的話這些 text placement 就不再以文字方塊處理）
    sn_done = False
    if sn_in_cell:
        sn_done = _write_sn_into_table(slide, placements, sn_col, sn_row_start)

    for pl in placements:
        if pl.text is not None:
            if sn_done:
                continue  # 已寫進儲存格
            tb = slide.shapes.add_textbox(Cm(pl.x_cm), Cm(pl.y_cm), Cm(pl.w_cm), Cm(pl.h_cm))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = align_map.get(pl.align, PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = pl.text
            run.font.size = Pt(pl.font_pt)
            run.font.bold = pl.bold
        else:
            # per-placement crop 優先；否則用全域 crop
            eff_crop = pl.crop or crop
            src = _apply_crop(pl.path, eff_crop, tmp_dir) if (eff_crop and tmp_dir) else pl.path
            slide.shapes.add_picture(
                str(src), Cm(pl.x_cm), Cm(pl.y_cm),
                width=Cm(pl.w_cm), height=Cm(pl.h_cm),
            )


def write_pages(
    slide_w_cm: float,
    slide_h_cm: float,
    pages: list[list[Placement]],
    out_path: Path,
    template: Path | None = None,
    sn_in_cell: bool = False,
    sn_col: int = 0,
    sn_row_start: int = 1,
    crop: dict | None = None,
) -> Path:
    """Write a multi-page pptx. Each page = a list of Placements.

    With a template: every page is a deep-copy of the template's first slide.
    Without a template: every page is a fresh blank slide.
    sn_in_cell=True 時 SN 文字寫進範本表格的儲存格 (而非浮動文字方塊)。
    """
    if not pages:
        raise ValueError("pages 不可為空")

    if template:
        prs = Presentation(str(template))
        if len(prs.slides) == 0:
            blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            base = prs.slides.add_slide(blank)
        else:
            base = prs.slides[0]
        # IMPORTANT: duplicate FIRST (while base is still pristine), THEN add images.
        # If we mutated base first, the duplicates would inherit those mutations.
        page_slides = [base]
        for _ in range(len(pages) - 1):
            page_slides.append(_duplicate_slide(prs, base))
        for slide, page_pls in zip(page_slides, pages):
            _add_placements_to_slide(slide, page_pls, sn_in_cell, sn_col, sn_row_start,
                                     crop=crop, tmp_dir=out_path.parent / "_crops")
    else:
        prs = Presentation()
        prs.slide_width = Cm(slide_w_cm)
        prs.slide_height = Cm(slide_h_cm)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        for page_pls in pages:
            slide = prs.slides.add_slide(blank)
            _add_placements_to_slide(slide, page_pls, sn_in_cell, sn_col, sn_row_start,
                                     crop=crop, tmp_dir=out_path.parent / "_crops")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _composite_cell_image(img_path: Path, cw_cm: float, rh_cm: float, fill: float,
                          crop: dict | None, tmp_dir: Path, fit: str = "contain") -> Path:
    """把照片合成到「儲存格比例的畫布」上 → 此合成圖填滿儲存格 → 照片跟著儲存格走（不漂移）。
    fit：contain(完整貼邊，等比留白) / cover(裁切填滿，等比裁切) / fill(拉伸變形填滿)。"""
    from .xlsx_writer import _apply_crop
    tmp_dir.mkdir(parents=True, exist_ok=True)
    src = _apply_crop(img_path, crop, tmp_dir) if crop else img_path
    DPI = 110
    W = max(8, int(cw_cm / 2.54 * DPI))
    H = max(8, int(rh_cm / 2.54 * DPI))
    im = Image.open(src).convert("RGB")
    if fit == "fill":
        canvas = im.resize((W, H))
    elif fit == "cover":
        # 等比放大到覆蓋整格，置中裁切
        scale = max(W / im.width, H / im.height)
        rw, rh = max(1, int(im.width * scale)), max(1, int(im.height * scale))
        im2 = im.resize((rw, rh))
        left, top = (rw - W) // 2, (rh - H) // 2
        canvas = im2.crop((left, top, left + W, top + H))
    else:  # contain
        canvas = Image.new("RGB", (W, H), "white")
        th = H * fill
        tw = th * im.width / im.height
        mw = W * fill
        if tw > mw:
            tw = mw
            th = tw * im.height / im.width
        im2 = im.resize((max(1, int(tw)), max(1, int(th))))
        canvas.paste(im2, ((W - int(tw)) // 2, (H - int(th)) // 2))
    out = tmp_dir / f"cell_{img_path.stem}_{fit}_{W}x{H}.png"
    canvas.save(out)
    return out


def _set_cell_picture_fill(slide, cell, image_path: Path) -> None:
    """把表格儲存格的填滿設成圖片（blipFill），圖片即成為儲存格內容、與格子完全貼齊。"""
    from pptx.oxml.ns import qn
    _img_part, rId = slide.part.get_or_add_image_part(str(image_path))
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    blipFill = tcPr.makeelement(qn("a:blipFill"), {})
    blipFill.append(tcPr.makeelement(qn("a:blip"), {qn("r:embed"): rId}))
    stretch = tcPr.makeelement(qn("a:stretch"), {})
    stretch.append(tcPr.makeelement(qn("a:fillRect"), {}))
    blipFill.append(stretch)
    ext = tcPr.find(qn("a:extLst"))
    if ext is not None:
        ext.addprevious(blipFill)
    else:
        tcPr.append(blipFill)


def write_sn_cell_pages(template: Path, out_path: Path, pages: list[dict],
                        fill: float = 0.9, fit: str = "contain") -> Path:
    """依範本 SN（PPT）：把 SN 文字與照片直接填進「表格儲存格」。
    照片成為儲存格內容（cell fill）→ 跟著儲存格走，任何檢視器都對齊（浮動圖會因列高渲染差異漂移）。
    pages: 每頁 { "sn": [{row,col,text,font_pt,bold}], "img": [{row,col,path,crop}] }（row/col 0-based）。
    """
    from pptx.util import Pt
    from pptx.util import Emu as _Emu
    if not pages:
        raise ValueError("pages 不可為空")
    prs = Presentation(str(template))
    base = prs.slides[0]
    page_slides = [base]
    for _ in range(len(pages) - 1):
        page_slides.append(_duplicate_slide(prs, base))
    tmp_dir = out_path.parent / "_cellimg"
    for slide, page in zip(page_slides, pages):
        table = _find_table(slide)
        if table is None:
            continue
        nrows, ncols = len(table.rows), len(table.columns)
        for sc in page.get("sn", []):
            r, c = int(sc["row"]), int(sc["col"])
            if not (0 <= r < nrows and 0 <= c < ncols):
                continue
            cell = table.cell(r, c)
            cell.text = str(sc.get("text", ""))
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(float(sc.get("font_pt", 14)))
                    run.font.bold = bool(sc.get("bold", True))
        for ic in page.get("img", []):
            r, c = int(ic["row"]), int(ic["col"])
            if not (0 <= r < nrows and 0 <= c < ncols) or not ic.get("path"):
                continue
            cw_cm = _Emu(table.columns[c].width).cm
            rh_cm = _Emu(table.rows[r].height).cm
            comp = _composite_cell_image(Path(ic["path"]), cw_cm, rh_cm, fill, ic.get("crop"), tmp_dir, fit)
            _set_cell_picture_fill(slide, table.cell(r, c), comp)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def write_placements(
    slide_w_cm: float,
    slide_h_cm: float,
    placements: list[Placement],
    out_path: Path,
    template: Path | None = None,
) -> Path:
    """Single-page convenience wrapper around write_pages."""
    return write_pages(slide_w_cm, slide_h_cm, [placements], out_path, template)


def _image_aspect(path: Path) -> float:
    """height / width"""
    with Image.open(path) as im:
        w, h = im.size
    return h / w if w else 1.0


def placements_from_config(cfg: Config, grouped: GroupedImages) -> list[Placement]:
    """Compute placements: each image keeps its natural aspect.

    cfg.grid.cell.w_cm = each image width
    cfg.grid.origin = first image top-left
    cfg.grid.gap = horizontal/vertical spacing between images
    Row height for row r = cell.w_cm * aspect(first non-null image of row r)
    """
    g = cfg.grid
    placements: list[Placement] = []
    cur_y = g.origin.y_cm

    for _group, row in grouped.rows:
        anchor = next((p for p in row if p is not None), None)
        if anchor is None:
            continue
        row_h = g.cell.w_cm * _image_aspect(anchor)
        for ci, p in enumerate(row):
            if p is None:
                continue
            x = g.origin.x_cm + ci * (g.cell.w_cm + g.gap.x_cm)
            placements.append(Placement(
                path=p, x_cm=x, y_cm=cur_y, w_cm=g.cell.w_cm, h_cm=row_h,
            ))
        cur_y += row_h + g.gap.y_cm
    return placements


def write_pptx(cfg: Config, grouped: GroupedImages) -> Path:
    """Backwards-compatible entry: computes placements from cfg, then writes."""
    placements = placements_from_config(cfg, grouped)
    return write_placements(
        cfg.slide.width_cm, cfg.slide.height_cm, placements, cfg.output.path, cfg.output.template,
    )
