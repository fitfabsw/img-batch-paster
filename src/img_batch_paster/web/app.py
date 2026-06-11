from __future__ import annotations

import io
import re
import shutil
import tempfile
import uuid
import json
from pathlib import Path


def _natural_key(name: str):
    """自然排序 key：把檔名數字段當數字比，讓 QN2 排在 QN10 前面。"""
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", name)]

import click
from flask import Flask, jsonify, request, send_file, send_from_directory
from PIL import Image

from ..grouper import scan_folder
from ..pptx_writer import Placement, write_pages, write_placements
from ..xlsx_writer import CellPlacement, write_xlsx
from ..keynote_export import convert_key_to_pptx, convert_pptx_to_key
from .template_render import render_first_slide, slide_size_cm, prewarm_libreoffice

STATIC_DIR = Path(__file__).parent / "static"
UPLOAD_DIR = Path(tempfile.gettempdir()) / "img-batch-paster-uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

app = Flask(__name__, static_folder=str(STATIC_DIR), static_url_path="/static")
app.config["MAX_CONTENT_LENGTH"] = 2 * 1024 * 1024 * 1024  # 2 GB per request


def _ws_dir(ws: str) -> Path | None:
    if not ws or "/" in ws or ".." in ws:
        return None
    p = UPLOAD_DIR / ws
    return p if p.is_dir() else None


@app.get("/")
def index():
    # no-store：HTML 內含整個 no-build app，禁止瀏覽器快取才不會卡在舊版（/static 的
    # React/Babel 大檔仍照常快取，不受影響）。免去每次改版都要手動硬重整。
    resp = send_from_directory(STATIC_DIR, "index.html")
    resp.headers["Cache-Control"] = "no-store"
    return resp


@app.post("/api/scan")
def api_scan():
    data = request.get_json(force=True)
    folder = Path(data["folder"]).expanduser().resolve()
    extensions = data.get("extensions", [".png", ".jpg", ".jpeg"])
    exts = {e.lower() for e in extensions}

    if not folder.is_dir():
        return jsonify({"error": f"資料夾不存在: {folder}"}), 400

    files = []
    for entry in sorted(folder.iterdir(), key=lambda p: _natural_key(p.name)):
        if not entry.is_file() or entry.suffix.lower() not in exts:
            continue
        try:
            with Image.open(entry) as im:
                w, h = im.size
        except Exception:
            w, h = 0, 0
        files.append({"path": str(entry), "name": entry.name, "w": w, "h": h})

    return jsonify({"folder": str(folder), "files": files})


@app.get("/api/thumb")
def api_thumb():
    path = request.args.get("path")
    if not path:
        return "missing path", 400
    p = Path(path)
    if not p.is_file():
        return "not found", 404
    size = int(request.args.get("size", 240))
    img = Image.open(p)
    # optional crop=l,t,r,b in 0..1
    crop_q = request.args.get("crop")
    if crop_q:
        try:
            l, t, r, b = (max(0.0, min(1.0, float(x))) for x in crop_q.split(","))
            if r > l and b > t and not (l == 0 and t == 0 and r == 1 and b == 1):
                w, h = img.size
                img = img.crop((int(l * w), int(t * h), int(r * w), int(b * h)))
        except Exception:
            pass
    img.thumbnail((size, size))
    buf = io.BytesIO()
    fmt = "PNG" if img.mode in ("RGBA", "P") else "JPEG"
    img.convert("RGB" if fmt == "JPEG" else img.mode).save(buf, fmt)
    buf.seek(0)
    return send_file(buf, mimetype=f"image/{fmt.lower()}")


@app.get("/api/version")
def api_version():
    import subprocess
    from .. import __version__
    project_root = Path(__file__).resolve().parents[3]
    def _git(*args):
        try:
            return subprocess.check_output(
                ["git", *args], cwd=project_root, stderr=subprocess.DEVNULL, text=True
            ).strip()
        except Exception:
            return ""
    return jsonify({
        "version": __version__,
        "branch": _git("rev-parse", "--abbrev-ref", "HEAD") or "?",
        "commit": _git("rev-parse", "--short", "HEAD") or "?",
        "date": _git("log", "-1", "--format=%cd", "--date=format:%Y-%m-%d %H:%M") or "?",
    })


@app.post("/api/workspace")
def api_workspace():
    """Create a fresh per-browser workspace under /tmp/img-batch-paster-uploads/."""
    ws = uuid.uuid4().hex[:12]
    (UPLOAD_DIR / ws / "images").mkdir(parents=True, exist_ok=True)
    return jsonify({"workspace": ws})


@app.post("/api/upload/template")
def api_upload_template():
    import sys as _sys
    ws = request.form.get("workspace", "")
    base = _ws_dir(ws)
    if not base:
        return jsonify({"error": "invalid workspace"}), 400
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": "no file"}), 400

    ext = Path(f.filename).suffix.lower()
    dest_pptx = base / "template.pptx"
    dest_xlsx = base / "template.xlsx"

    if ext == ".pptx":
        f.save(str(dest_pptx))
        return jsonify({"path": str(dest_pptx), "mode": "slides"})

    if ext == ".xlsx":
        f.save(str(dest_xlsx))
        return jsonify({"path": str(dest_xlsx), "mode": "excel"})

    if ext == ".key":
        if _sys.platform != "darwin":
            return jsonify({"error": ".key 範本僅在 macOS server 可轉檔；請改傳 .pptx"}), 501
        tmp_key = base / "uploaded.key"
        f.save(str(tmp_key))
        try:
            convert_key_to_pptx(tmp_key, dest_pptx)
        except Exception as e:
            return jsonify({"error": f"Keynote → pptx 轉檔失敗: {e}"}), 500
        finally:
            if tmp_key.exists():
                tmp_key.unlink()
        return jsonify({"path": str(dest_pptx), "mode": "slides"})

    return jsonify({"error": f"不支援的副檔名: {ext}（請用 .pptx / .key / .xlsx）"}), 400


DEFAULT_TEMPLATE = Path(__file__).resolve().parent.parent / "templates" / "default.pptx"


@app.post("/api/upload/images")
def api_upload_images():
    ws = request.form.get("workspace", "")
    base = _ws_dir(ws)
    if not base:
        return jsonify({"error": "invalid workspace"}), 400
    folder = base / "images"
    # Replace any previous uploads in this workspace
    if folder.exists():
        shutil.rmtree(folder)
    folder.mkdir(parents=True, exist_ok=True)

    files = request.files.getlist("files")
    saved = 0
    for f in files:
        if not f.filename:
            continue
        # 只保留檔名，剝掉路徑成分；保留 unicode（避免 secure_filename 過濾中文）
        name = Path(f.filename).name
        if not name:
            continue
        f.save(str(folder / name))
        saved += 1
    return jsonify({"folder": str(folder), "count": saved})


@app.post("/api/upload/source")
def api_upload_source():
    """Upload one source folder for SN-match mode. Workspace gets source-N/ subfolders."""
    ws = request.form.get("workspace", "")
    base = _ws_dir(ws)
    if not base:
        return jsonify({"error": "invalid workspace"}), 400
    try:
        sid = int(request.form.get("sourceId", "0"))
    except ValueError:
        return jsonify({"error": "invalid sourceId"}), 400
    folder = base / f"source-{sid}"
    if folder.exists():
        shutil.rmtree(folder)
    folder.mkdir(parents=True, exist_ok=True)

    files = request.files.getlist("files")
    saved_files = []
    for f in files:
        if not f.filename:
            continue
        name = Path(f.filename).name
        if not name:
            continue
        target = folder / name
        f.save(str(target))
        try:
            with Image.open(target) as im:
                w, h = im.size
        except Exception:
            w, h = 0, 0
        saved_files.append({"name": name, "path": str(target), "w": w, "h": h})
    return jsonify({"folder": str(folder), "count": len(saved_files), "files": saved_files})


@app.post("/api/template/sn-list")
def api_template_sn_list():
    """Read SN list from an Excel template's column (e.g. column A from row 2)."""
    from openpyxl import load_workbook
    data = request.get_json(force=True)
    path = Path(data["path"]).expanduser().resolve()
    if not path.is_file():
        return jsonify({"error": f"檔案不存在: {path}"}), 400
    col_letter = (data.get("col") or "A").upper()
    try:
        row_start = max(1, int(data.get("rowStart", 2)))
    except (TypeError, ValueError):
        row_start = 2
    try:
        wb = load_workbook(str(path), data_only=True)
    except Exception as e:
        return jsonify({"error": f"無法開啟 Excel: {e}"}), 500
    ws = wb.active
    from openpyxl.utils import column_index_from_string
    try:
        col_idx = column_index_from_string(col_letter)
    except ValueError:
        return jsonify({"error": f"無效欄位字母: {col_letter}"}), 400
    sns = []
    # 連續讀直到遇到空 cell；同時記錄 row 號方便寫入時對應
    max_scan = (ws.max_row or row_start) + 200  # 容錯避免無限
    r = row_start
    while r <= max_scan:
        v = ws.cell(row=r, column=col_idx).value
        if v is None or (isinstance(v, str) and not v.strip()):
            # 允許中間出現 1-2 個空白後仍繼續嗎? 先採嚴格：遇空即停
            break
        sns.append({"sn": str(v).strip(), "row": r})
        r += 1
    return jsonify({"col": col_letter, "rowStart": row_start, "count": len(sns), "sns": sns})


@app.get("/api/download/<ws>/<path:filename>")
def api_download(ws, filename):
    base = _ws_dir(ws)
    if not base:
        return "not found", 404
    p = base / filename
    if not p.is_file() or ".." in filename:
        return "not found", 404
    return send_file(str(p), as_attachment=True, download_name=Path(filename).name)


@app.post("/api/pick")
def api_pick():
    """用 macOS osascript 開原生檔案/資料夾選擇器。"""
    import subprocess
    import sys
    if sys.platform != "darwin":
        return jsonify({"error": "檔案選擇器僅在 macOS 主機可用；請手動輸入路徑"}), 501
    data = request.get_json(force=True) or {}
    kind = data.get("kind", "folder")  # "folder" | "file"
    prompt = data.get("prompt", "請選擇")

    # 起始位置：以欄位現有路徑為基準（找最近存在的祖先目錄）
    default = data.get("default", "")
    default_clause = ""
    if default:
        start = Path(default).expanduser()
        if not start.is_dir():
            start = start.parent
        while start != start.parent and not start.is_dir():
            start = start.parent
        if start.is_dir():
            default_clause = f' default location POSIX file "{start}"'

    if kind == "folder":
        script = f'POSIX path of (choose folder with prompt "{prompt}"{default_clause})'
    else:
        ext_filter = data.get("extensions")
        if ext_filter:
            of = "{" + ", ".join(f'"{e.lstrip(".")}"' for e in ext_filter) + "}"
            script = f'POSIX path of (choose file with prompt "{prompt}" of type {of}{default_clause})'
        else:
            script = f'POSIX path of (choose file with prompt "{prompt}"{default_clause})'

    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=300,
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    if result.returncode != 0:
        # 使用者取消會回 1；其它錯誤同樣回空字串
        return jsonify({"path": None, "cancelled": True})
    path = result.stdout.strip().rstrip("/")
    return jsonify({"path": path})


@app.post("/api/template/load")
def api_template_load():
    data = request.get_json(force=True)
    path = Path(data["path"]).expanduser().resolve()
    if not path.is_file():
        return jsonify({"error": f"檔案不存在: {path}"}), 400
    try:
        w_cm, h_cm = slide_size_cm(path)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    # 渲染預覽 (需要 LibreOffice)；失敗就無背景，但仍回傳尺寸供 layout 使用
    preview_url = None
    try:
        png = render_first_slide(path)
        preview_url = f"/api/template/preview?key={png.stem}"
    except Exception:
        pass
    return jsonify({
        "path": str(path),
        "width_cm": w_cm,
        "height_cm": h_cm,
        "preview_url": preview_url,
    })


@app.post("/api/template/table-info")
def api_template_table_info():
    """讀取 .pptx 範本中所有表格的幾何（列/欄/位置/各欄寬/各列高），供自動貼圖用。
    .key 範本在上傳時已轉成 .pptx，故此處一律吃 .pptx 路徑。"""
    from pptx import Presentation
    from pptx.util import Emu
    data = request.get_json(force=True)
    path = Path(data["path"]).expanduser().resolve()
    if not path.is_file():
        return jsonify({"error": f"檔案不存在: {path}"}), 400
    if path.suffix.lower() != ".pptx":
        return jsonify({"error": f"只能讀 .pptx 表格（.key 上傳時會自動轉檔），目前: {path.suffix}"}), 400
    try:
        prs = Presentation(str(path))
    except Exception as e:
        return jsonify({"error": f"無法讀取簡報: {e}"}), 500
    sw, sh = Emu(prs.slide_width).cm, Emu(prs.slide_height).cm
    tables = []
    for si, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if not getattr(shp, "has_table", False):
                continue
            t = shp.table
            # 有文字的儲存格（1-based），供前端自動偵測表頭列 / 標籤欄
            cells = []
            for ri in range(len(t.rows)):
                for ci in range(len(t.columns)):
                    try:
                        txt = t.cell(ri, ci).text.strip()
                    except Exception:
                        txt = ""
                    if txt:
                        cells.append({"r": ri + 1, "c": ci + 1, "text": txt})
            tables.append({
                "slide": si,
                "name": shp.name,
                "rows": len(t.rows),
                "cols": len(t.columns),
                "left_cm": round(Emu(shp.left).cm, 2),
                "top_cm": round(Emu(shp.top).cm, 2),
                "width_cm": round(Emu(shp.width).cm, 2),
                "height_cm": round(Emu(shp.height).cm, 2),
                "col_widths_cm": [round(Emu(c.width).cm, 2) for c in t.columns],
                "row_heights_cm": [round(Emu(r.height).cm, 2) for r in t.rows],
                "cells": cells,
            })
    return jsonify({
        "slide": {"width_cm": round(sw, 2), "height_cm": round(sh, 2)},
        "tables": tables,
    })


@app.post("/api/template/excel-grid")
def api_template_excel_grid():
    """Return cell-grid structure of an .xlsx template for frontend preview."""
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    data = request.get_json(force=True)
    path = Path(data["path"]).expanduser().resolve()
    if not path.is_file():
        return jsonify({"error": f"檔案不存在: {path}"}), 400
    try:
        wb = load_workbook(str(path))
    except Exception as e:
        return jsonify({"error": f"無法讀取: {e}"}), 500
    ws = wb.active
    max_col = max(ws.max_column or 1, 12)
    max_row = max(ws.max_row or 1, 20)
    DEFAULT_W = 8.43
    DEFAULT_H = 15.0

    # 收集 range-based 欄寬：openpyxl 可能用 <col min=A max=B width=W> 一次涵蓋多欄
    col_widths: dict[int, float] = {}
    for cd in ws.column_dimensions.values():
        if cd.min is None or cd.max is None or cd.width is None:
            continue
        for i in range(cd.min, cd.max + 1):
            col_widths[i] = cd.width

    # 偵測預設字體大小調整 MDW (Max Digit Width)
    try:
        sz = float(wb._fonts[0].sz or 11)
    except Exception:
        sz = 11.0
    mdw = 7.0 if sz <= 11 else 7.0 + (sz - 11) * 0.85

    cols = []
    for c in range(1, max_col + 4):
        letter = get_column_letter(c)
        w = col_widths.get(c, DEFAULT_W)
        cols.append({"letter": letter, "w_px": w * mdw + 5})
    rows = []
    for r in range(1, max_row + 4):
        rd = ws.row_dimensions.get(r)
        h = rd.height if rd is not None and rd.height is not None else DEFAULT_H
        rows.append({"r": r, "h_px": h * 4 / 3})

    cells = []
    for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
        for cell in row:
            if cell.value is None:
                continue
            cells.append({
                "r": cell.row, "c": cell.column,
                "text": str(cell.value),
                "font_pt": float(cell.font.size or 11),
                "bold": bool(cell.font.bold),
                "h_align": cell.alignment.horizontal or "left",
                "v_align": cell.alignment.vertical or "bottom",
            })

    # Borders (簡化：只要有設定就算有；不分四邊)
    borders = []
    for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
        for cell in row:
            b = cell.border
            sides = {
                "top": bool(b.top and b.top.style),
                "right": bool(b.right and b.right.style),
                "bottom": bool(b.bottom and b.bottom.style),
                "left": bool(b.left and b.left.style),
            }
            if any(sides.values()):
                borders.append({"r": cell.row, "c": cell.column, **sides})

    return jsonify({
        "sheet": ws.title,
        "cols": cols,
        "rows": rows,
        "cells": cells,
        "borders": borders,
    })


@app.get("/api/template/preview")
def api_template_preview():
    from .template_render import CACHE_DIR
    key = request.args.get("key", "")
    p = CACHE_DIR / f"{key}.png"
    if not p.is_file():
        return "not found", 404
    return send_file(str(p), mimetype="image/png")


# --- 設定檔配方 (Recipe) CRUD ---
# 存 ~/.img-batch-paster/configs/<name>.ibp (zip 內含 manifest.json + template.<ext>)
# 向下相容：舊版 .json 純文字檔也讀得到 (但寫入一律 .ibp)
CONFIG_DIR = Path.home() / ".img-batch-paster" / "configs"
import re as _re
import zipfile as _zipfile


def _safe_config_name(name: str) -> str | None:
    """允許中英數、空白、底線、減號、點。其他字元一律拒絕避免 path traversal。"""
    name = (name or "").strip()
    if not name or len(name) > 80:
        return None
    if not _re.match(r"^[\w \-.一-鿿]+$", name):
        return None
    return name


def _config_file(name: str) -> Path | None:
    """找出 <name>.ibp 或舊版 <name>.json，回傳實際存在的那個 path；都沒有回 None。"""
    p_ibp = CONFIG_DIR / f"{name}.ibp"
    if p_ibp.is_file():
        return p_ibp
    p_json = CONFIG_DIR / f"{name}.json"
    if p_json.is_file():
        return p_json
    return None


def _read_recipe(p: Path) -> dict | None:
    """從 .ibp 或 .json 讀 recipe；失敗回 None。"""
    try:
        if p.suffix.lower() == ".ibp":
            with _zipfile.ZipFile(p, "r") as zf:
                with zf.open("manifest.json") as f:
                    return json.loads(f.read().decode("utf-8"))
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


def _ibp_template_name(p: Path) -> str | None:
    """讀 .ibp 內含的 template 檔名 (e.g. template.xlsx)；沒有則 None。"""
    if p.suffix.lower() != ".ibp":
        return None
    try:
        with _zipfile.ZipFile(p, "r") as zf:
            for n in zf.namelist():
                if n.startswith("template."):
                    return n
    except Exception:
        pass
    return None


@app.get("/api/configs")
def api_configs_list():
    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    items = []
    seen = set()
    # 先讀 .ibp，再 fallback .json (避免重名情況下 .ibp 優先)
    for pattern in ("*.ibp", "*.json"):
        for f in CONFIG_DIR.glob(pattern):
            if f.stem in seen:
                continue
            seen.add(f.stem)
            recipe = _read_recipe(f) or {}
            items.append({
                "name": f.stem,
                "format": f.suffix.lstrip("."),
                "hasTemplate": _ibp_template_name(f) is not None,
                "meta": recipe.get("_meta", {}),
                "mtime": f.stat().st_mtime,
            })
    items.sort(key=lambda x: x["mtime"], reverse=True)
    return jsonify({"configs": items})


@app.post("/api/configs")
def api_configs_save():
    """儲存配置 — 寫成 .ibp (zip)，可附帶範本檔。
    請求 JSON: { name, recipe, templatePath? }
    """
    data = request.get_json(force=True)
    name = _safe_config_name(data.get("name", ""))
    if not name:
        return jsonify({"error": "無效的配置名稱"}), 400
    recipe = data.get("recipe") or {}
    if not isinstance(recipe, dict):
        return jsonify({"error": "recipe 必須是 object"}), 400
    template_path_str = data.get("templatePath") or ""
    tpl_path = Path(template_path_str).expanduser().resolve() if template_path_str else None

    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    out = CONFIG_DIR / f"{name}.ibp"
    # 刪掉舊的 .json（若有）以免兩份混淆
    old_json = CONFIG_DIR / f"{name}.json"
    if old_json.is_file():
        old_json.unlink()

    with _zipfile.ZipFile(out, "w", _zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("manifest.json", json.dumps(recipe, ensure_ascii=False, indent=2))
        if tpl_path and tpl_path.is_file():
            inner_name = f"template{tpl_path.suffix}"
            zf.write(str(tpl_path), arcname=inner_name)
    return jsonify({"name": name, "path": str(out), "hasTemplate": bool(tpl_path and tpl_path.is_file())})


@app.get("/api/configs/<name>")
def api_configs_get(name):
    """載入配置 — 回傳 recipe + (若有範本) 解壓到 workspace 的 template path。"""
    safe = _safe_config_name(name)
    if not safe:
        return jsonify({"error": "無效的配置名稱"}), 400
    p = _config_file(safe)
    if p is None:
        return jsonify({"error": f"配置不存在: {safe}"}), 404
    recipe = _read_recipe(p)
    if recipe is None:
        return jsonify({"error": "無法讀取配置內容"}), 500

    resp = dict(recipe)
    resp["_loaded"] = {"name": safe, "format": p.suffix.lstrip(".")}

    # 如果 .ibp 含範本，解壓到 workspace
    ws = request.args.get("workspace", "")
    base = _ws_dir(ws) if ws else None
    if base and p.suffix.lower() == ".ibp":
        inner = _ibp_template_name(p)
        if inner:
            ext = Path(inner).suffix
            dest = base / f"template{ext}"
            with _zipfile.ZipFile(p, "r") as zf:
                with zf.open(inner) as src, open(dest, "wb") as dst:
                    dst.write(src.read())
            resp["_loaded"]["templatePath"] = str(dest)
            resp["_loaded"]["templateName"] = recipe.get("_meta", {}).get("templateHint") or inner
    return jsonify(resp)


@app.delete("/api/configs/<name>")
def api_configs_delete(name):
    safe = _safe_config_name(name)
    if not safe:
        return jsonify({"error": "無效的配置名稱"}), 400
    p = _config_file(safe)
    if p is None:
        return jsonify({"error": "配置不存在"}), 404
    p.unlink()
    return jsonify({"ok": True})


@app.get("/api/configs/<name>/download")
def api_configs_download(name):
    """下載 .ibp / .json 配置檔給使用者帶走。"""
    safe = _safe_config_name(name)
    if not safe:
        return "invalid name", 400
    p = _config_file(safe)
    if p is None:
        return "not found", 404
    return send_file(str(p), as_attachment=True, download_name=p.name)


@app.post("/api/configs/import")
def api_configs_import():
    """匯入使用者上傳的 .ibp 或 .json 配置檔；存到 CONFIG_DIR 後可在下拉看到。"""
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": "no file"}), 400
    src_name = Path(f.filename).name
    stem = Path(src_name).stem
    ext = Path(src_name).suffix.lower()
    if ext not in (".ibp", ".json"):
        return jsonify({"error": f"不支援的副檔名: {ext}"}), 400
    safe = _safe_config_name(stem)
    if not safe:
        return jsonify({"error": "檔名含不允許字元"}), 400
    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    out = CONFIG_DIR / f"{safe}{ext}"
    f.save(str(out))
    # 驗證內容
    recipe = _read_recipe(out)
    if recipe is None:
        out.unlink()
        return jsonify({"error": "檔案內容無效"}), 400
    return jsonify({"name": safe, "format": ext.lstrip("."), "hasTemplate": _ibp_template_name(out) is not None})


@app.post("/api/export")
def api_export():
    data = request.get_json(force=True)
    slide = data.get("slide", {"width_cm": 25.4, "height_cm": 14.29})

    ws = data.get("workspace") or data["output"].get("workspace")
    requested = data["output"]["path"]
    base = _ws_dir(ws) if ws else None
    if base:
        out_name = Path(requested).name or "out.pptx"
        out_path = base / out_name
    else:
        out_path = Path(requested).expanduser().resolve()

    template = data["output"].get("template")
    template_path = Path(template).expanduser().resolve() if template else None

    # Excel 分支：副檔名 .xlsx
    if out_path.suffix.lower() == ".xlsx":
        xl_placements = [
            CellPlacement(
                path=Path(p["path"]) if p.get("path") else None,
                row=int(p["row"]), col=int(p["col"]),
                span_cols=int(p.get("span_cols", 1)),
                span_rows=int(p.get("span_rows", 1)),
                text=p.get("text"),
                font_pt=float(p.get("font_pt", 12)),
                crop=p.get("crop"),
            )
            for p in data.get("cells", [])
        ]
        if not xl_placements:
            return jsonify({"error": "沒有可匯出的儲存格"}), 400
        embed_in_cell = bool(data.get("embedInCell"))
        img_fit = data.get("imgFit", "cover")
        contain_inset = float(data.get("containInset", 0.05))
        crop = data.get("crop")
        write_xlsx(xl_placements, out_path, template_path,
                   embed_in_cell=embed_in_cell, img_fit=img_fit,
                   contain_inset=contain_inset, crop=crop)
        resp = {"output": str(out_path)}
        if base:
            resp["download_url"] = f"/api/download/{ws}/{out_path.name}"
        return jsonify(resp)

    def _to_pl(p):
        return Placement(
            path=Path(p["path"]) if p.get("path") else None,
            x_cm=float(p["x_cm"]), y_cm=float(p["y_cm"]),
            w_cm=float(p["w_cm"]), h_cm=float(p["h_cm"]),
            text=p.get("text"),
            font_pt=float(p.get("font_pt", 18)),
            bold=bool(p.get("bold", True)),
            align=p.get("align", "center"),
            row_idx=p.get("row_idx"),
            crop=p.get("crop"),
        )

    if "pages" in data:
        pages = [[_to_pl(p) for p in page] for page in data["pages"]]
    else:
        pages = [[_to_pl(p) for p in data.get("placements", [])]]

    if not pages or all(len(p) == 0 for p in pages):
        return jsonify({"error": "沒有可匯出的圖片"}), 400

    sn_in_cell = bool(data.get("snInCell"))
    sn_col = int(data.get("snCol", 0))
    sn_row_start = int(data.get("snRowStart", 1))
    crop = data.get("crop")

    # 副檔名 .key → 先輸出 .pptx，再經由 Keynote 轉成 .key
    want_key = out_path.suffix.lower() == ".key"
    pptx_out = out_path.with_suffix(".pptx") if want_key else out_path
    write_pages(
        float(slide["width_cm"]), float(slide["height_cm"]),
        pages, pptx_out, template_path,
        sn_in_cell=sn_in_cell, sn_col=sn_col, sn_row_start=sn_row_start,
        crop=crop,
    )

    if want_key:
        try:
            convert_pptx_to_key(pptx_out, out_path)
        except Exception as e:
            return jsonify({
                "error": f".key 轉檔失敗，但 .pptx 已輸出至 {pptx_out}: {e}",
                "output": str(pptx_out), "pages": len(pages),
            }), 500
        # 保留中介 .pptx 供 debug；要刪可改為 pptx_out.unlink(missing_ok=True)

    resp = {"output": str(out_path), "pages": len(pages)}
    if base:
        resp["download_url"] = f"/api/download/{ws}/{out_path.name}"
    return jsonify(resp)


@click.command()
@click.option("--host", default="127.0.0.1")
@click.option("--port", default=5050, type=int)
@click.option("--desktop", is_flag=True, help="以 pywebview 桌面視窗開啟")
def run(host: str, port: int, desktop: bool) -> None:
    """啟動 Web 預覽伺服器。"""
    # 背景暖機 LibreOffice profile，讓使用者第一次上傳範本時不用等冷啟動
    prewarm_libreoffice(DEFAULT_TEMPLATE)
    if desktop:
        try:
            import webview  # type: ignore
        except ImportError:
            raise SystemExit("請先安裝桌面依賴: pip install -e '.[desktop]'")
        import threading

        threading.Thread(
            target=lambda: app.run(host=host, port=port, debug=False, use_reloader=False),
            daemon=True,
        ).start()
        webview.create_window("img-batch-paster", f"http://{host}:{port}", width=1200, height=820)
        webview.start()
    else:
        click.echo(f" * 開啟 http://{host}:{port}")
        app.run(host=host, port=port, debug=True, use_reloader=False)


if __name__ == "__main__":
    run()
